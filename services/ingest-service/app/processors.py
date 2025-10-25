import os
import fitz  # PyMuPDF
import docx
import pptx
import whisper
import openai
import weaviate
from pathlib import Path
from typing import List, Dict, Any
import tiktoken
import uuid
import logging

logger = logging.getLogger(__name__)

class BaseProcessor:
    """Base class for file processors"""
    
    def __init__(self):
        self.openai_client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.weaviate_client = weaviate.Client(url=os.getenv("WEAVIATE_URL", "http://localhost:8080"))
        self.encoding = tiktoken.get_encoding("cl100k_base")
    
    def chunk_text(self, text: str, max_tokens: int = 500, overlap: int = 50) -> List[str]:
        """Chunk text with overlap"""
        tokens = self.encoding.encode(text)
        chunks = []
        
        start = 0
        while start < len(tokens):
            end = min(start + max_tokens, len(tokens))
            chunk_tokens = tokens[start:end]
            chunk_text = self.encoding.decode(chunk_tokens)
            chunks.append(chunk_text)
            
            if end >= len(tokens):
                break
            
            start = end - overlap
        
        return chunks
    
    def create_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Create embeddings for text chunks"""
        try:
            response = self.openai_client.embeddings.create(
                model="text-embedding-ada-002",
                input=texts
            )
            return [embedding.embedding for embedding in response.data]
        except Exception as e:
            logger.error(f"Failed to create embeddings: {e}")
            raise
    
    def store_vectors(self, chunks: List[str], embeddings: List[List[float]], 
                     source_id: str, persona_id: str, org_id: str) -> int:
        """Store vectors in Weaviate"""
        try:
            # Ensure schema exists
            self.ensure_schema()
            
            stored_count = 0
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                chunk_id = str(uuid.uuid4())
                
                data_object = {
                    "text": chunk,
                    "source_id": source_id,
                    "persona_id": persona_id,
                    "org_id": org_id,
                    "chunk_index": i,
                    "chunk_id": chunk_id
                }
                
                self.weaviate_client.data_object.create(
                    data_object=data_object,
                    class_name="KnowledgeChunk",
                    vector=embedding
                )
                stored_count += 1
            
            return stored_count
            
        except Exception as e:
            logger.error(f"Failed to store vectors: {e}")
            raise
    
    def ensure_schema(self):
        """Ensure Weaviate schema exists"""
        try:
            # Check if class exists
            schema = self.weaviate_client.schema.get()
            class_names = [cls["class"] for cls in schema.get("classes", [])]
            
            if "KnowledgeChunk" not in class_names:
                # Create schema
                class_obj = {
                    "class": "KnowledgeChunk",
                    "description": "Knowledge chunks for RAG",
                    "properties": [
                        {
                            "name": "text",
                            "dataType": ["text"],
                            "description": "The text content of the chunk"
                        },
                        {
                            "name": "source_id",
                            "dataType": ["string"],
                            "description": "ID of the source document"
                        },
                        {
                            "name": "persona_id",
                            "dataType": ["string"],
                            "description": "ID of the associated persona"
                        },
                        {
                            "name": "org_id",
                            "dataType": ["string"],
                            "description": "ID of the organization"
                        },
                        {
                            "name": "chunk_index",
                            "dataType": ["int"],
                            "description": "Index of the chunk in the document"
                        },
                        {
                            "name": "chunk_id",
                            "dataType": ["string"],
                            "description": "Unique ID of the chunk"
                        }
                    ]
                }
                
                self.weaviate_client.schema.create_class(class_obj)
                
        except Exception as e:
            logger.error(f"Failed to ensure schema: {e}")
            raise

class DocumentProcessor(BaseProcessor):
    """Processor for document files (PDF, DOCX, PPTX, TXT)"""
    
    def process(self, file_path: str, source_id: str, persona_id: str, org_id: str) -> Dict[str, Any]:
        """Process document file"""
        file_path = Path(file_path)
        
        # Extract text based on file type
        if file_path.suffix.lower() == '.pdf':
            text = self.extract_pdf_text(file_path)
        elif file_path.suffix.lower() == '.docx':
            text = self.extract_docx_text(file_path)
        elif file_path.suffix.lower() == '.pptx':
            text = self.extract_pptx_text(file_path)
        elif file_path.suffix.lower() == '.txt':
            text = self.extract_txt_text(file_path)
        else:
            raise ValueError(f"Unsupported document type: {file_path.suffix}")
        
        # Chunk text
        chunks = self.chunk_text(text)
        
        # Create embeddings
        embeddings = self.create_embeddings(chunks)
        
        # Store in vector database
        chunk_count = self.store_vectors(chunks, embeddings, source_id, persona_id, org_id)
        
        return {"chunk_count": chunk_count}
    
    def extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF"""
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text
    
    def extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    
    def extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PPTX"""
        presentation = pptx.Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def extract_txt_text(self, file_path: Path) -> str:
        """Extract text from TXT"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

class AudioProcessor(BaseProcessor):
    """Processor for audio files"""
    
    def __init__(self):
        super().__init__()
        self.whisper_model = whisper.load_model("base")
    
    def process(self, file_path: str, source_id: str, persona_id: str, org_id: str) -> Dict[str, Any]:
        """Process audio file"""
        # Transcribe audio
        result = self.whisper_model.transcribe(file_path)
        text = result["text"]
        
        # Chunk text
        chunks = self.chunk_text(text)
        
        # Create embeddings
        embeddings = self.create_embeddings(chunks)
        
        # Store in vector database
        chunk_count = self.store_vectors(chunks, embeddings, source_id, persona_id, org_id)
        
        return {"chunk_count": chunk_count}

class VideoProcessor(BaseProcessor):
    """Processor for video files"""
    
    def __init__(self):
        super().__init__()
        self.whisper_model = whisper.load_model("base")
    
    def process(self, file_path: str, source_id: str, persona_id: str, org_id: str) -> Dict[str, Any]:
        """Process video file (extract audio and transcribe)"""
        # For now, use whisper directly on video file
        # In production, you might want to extract audio first using ffmpeg
        result = self.whisper_model.transcribe(file_path)
        text = result["text"]
        
        # Chunk text
        chunks = self.chunk_text(text)
        
        # Create embeddings
        embeddings = self.create_embeddings(chunks)
        
        # Store in vector database
        chunk_count = self.store_vectors(chunks, embeddings, source_id, persona_id, org_id)
        
        return {"chunk_count": chunk_count}